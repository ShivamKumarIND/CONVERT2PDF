"""
Convert TO PDF Tools
- JPG to PDF
- WORD to PDF
- POWERPOINT to PDF
- EXCEL to PDF
- HTML to PDF
"""
from pathlib import Path
from typing import BinaryIO, List
from PIL import Image
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import io

class ConvertToPDF:
    
    @staticmethod
    def image_to_pdf(image_files: List[BinaryIO], output_path: Path) -> Path:
        """
        Convert images (JPG, PNG, etc.) to PDF
        
        Args:
            image_files: List of image file objects
            output_path: Path to save PDF
            
        Returns:
            Path to created PDF
        """
        try:
            images = []
            for img_file in image_files:
                img = Image.open(img_file)
                # Convert to RGB if needed (PDF doesn't support RGBA)
                if img.mode in ('RGBA', 'LA', 'P'):
                    # Create white background
                    rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    rgb_img.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                    img = rgb_img
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                images.append(img)
            
            # Save as PDF
            if images:
                images[0].save(
                    output_path,
                    "PDF",
                    save_all=True,
                    append_images=images[1:] if len(images) > 1 else []
                )
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting image to PDF: {str(e)}")
    
    @staticmethod
    def word_to_pdf(docx_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert Word document to PDF
        
        Note: This is a simplified version. For production, use:
        - docx2pdf library (Windows only)
        - LibreOffice in headless mode (cross-platform)
        - Cloud conversion API
        
        Args:
            docx_file: Word file object
            output_path: Path to save PDF
            
        Returns:
            Path to PDF
        """
        try:
            # This is a placeholder implementation
            # Real implementation would require additional dependencies
            
            # Option 1: Use docx2pdf (Windows only)
            # from docx2pdf import convert
            # convert(docx_file, output_path)
            
            # Option 2: Use LibreOffice
            # subprocess.run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, docx_path])
            
            # For now, we'll create a basic PDF from Word content
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            
            doc = Document(docx_file)
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            for para in doc.paragraphs:
                if para.text.strip():
                    p = Paragraph(para.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            pdf_doc.build(story)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting Word to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_pdf(excel_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert Excel to PDF
        
        Args:
            excel_file: Excel file object
            output_path: Path to save PDF
            
        Returns:
            Path to PDF
        """
        try:
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            
            # Load workbook
            wb = load_workbook(excel_file, data_only=True)
            ws = wb.active
            
            # Get data
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else '' for cell in row])
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(str(output_path), pagesize=landscape(letter))
            
            # Create table
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            pdf_doc.build([table])
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting Excel to PDF: {str(e)}")
    
    @staticmethod
    def powerpoint_to_pdf(pptx_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert PowerPoint to PDF
        
        Args:
            pptx_file: PowerPoint file object
            output_path: Path to save PDF
            
        Returns:
            Path to PDF
        """
        try:
            from reportlab.lib.pagesizes import letter, landscape
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import inch
            
            # Load presentation
            prs = Presentation(pptx_file)
            
            # Create PDF
            c = canvas.Canvas(str(output_path), pagesize=landscape(letter))
            width, height = landscape(letter)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Add slide number
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, height - 50, f"Slide {slide_num}")
                
                y_position = height - 100
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        c.setFont("Helvetica", 12)
                        text = shape.text[:100]  # Limit text length
                        c.drawString(50, y_position, text)
                        y_position -= 20
                
                c.showPage()
            
            c.save()
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting PowerPoint to PDF: {str(e)}")
    
    @staticmethod
    def html_to_pdf(html_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert HTML to PDF
        
        Args:
            html_file: HTML file object
            output_path: Path to save PDF
            
        Returns:
            Path to PDF
        """
        try:
            from weasyprint import HTML, CSS
            
            # Read HTML content
            html_content = html_file.read().decode('utf-8')
            
            # Convert to PDF
            HTML(string=html_content).write_pdf(output_path)
            
            return output_path
            
        except Exception as e:
            # Fallback to basic conversion if weasyprint fails
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                
                c = canvas.Canvas(str(output_path), pagesize=letter)
                c.setFont("Helvetica", 12)
                
                html_file.seek(0)
                html_content = html_file.read().decode('utf-8')
                
                # Simple text extraction (remove HTML tags)
                import re
                text = re.sub('<[^<]+?>', '', html_content)
                
                y = 750
                for line in text.split('\n')[:40]:  # Limit lines
                    if line.strip():
                        c.drawString(50, y, line[:80])
                        y -= 15
                
                c.save()
                return output_path
                
            except Exception as e2:
                raise Exception(f"Error converting HTML to PDF: {str(e2)}")
