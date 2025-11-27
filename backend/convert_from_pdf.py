"""
Convert FROM PDF Tools
- PDF to JPG
- PDF to WORD
- PDF to POWERPOINT
- PDF to EXCEL
- PDF to PDF/A
"""
from pathlib import Path
from typing import BinaryIO, List
from PyPDF2 import PdfReader
import pikepdf
from PIL import Image
import io

class ConvertFromPDF:
    
    @staticmethod
    def pdf_to_images(pdf_file: BinaryIO, output_dir: Path, 
                      image_format: str = "jpg", dpi: int = 200) -> List[Path]:
        """
        Convert PDF pages to images
        
        Args:
            pdf_file: PDF file object
            output_dir: Directory to save images
            image_format: 'jpg' or 'png'
            dpi: Image resolution (default: 200)
            
        Returns:
            List of paths to created images
        """
        try:
            # Try using pdf2image if available
            try:
                from pdf2image import convert_from_bytes
                
                pdf_file.seek(0)
                pdf_bytes = pdf_file.read()
                
                images = convert_from_bytes(pdf_bytes, dpi=dpi)
                output_files = []
                
                for i, image in enumerate(images, 1):
                    output_path = output_dir / f"page_{i}.{image_format}"
                    
                    if image_format.lower() == 'jpg':
                        # Convert to RGB for JPEG
                        if image.mode != 'RGB':
                            image = image.convert('RGB')
                        image.save(output_path, 'JPEG', quality=95)
                    else:
                        image.save(output_path, 'PNG')
                    
                    output_files.append(output_path)
                
                return output_files
                
            except ImportError:
                # Fallback: This won't work well but provides a placeholder
                raise Exception("pdf2image not available. Install with: pip install pdf2image")
                
        except Exception as e:
            raise Exception(f"Error converting PDF to images: {str(e)}")
    
    @staticmethod
    def pdf_to_word(pdf_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert PDF to Word document
        
        Note: This is a basic implementation. For production quality:
        - Use Adobe PDF Services API
        - Use pdf2docx library
        - Use cloud conversion services
        
        Args:
            pdf_file: PDF file object
            output_path: Path to save DOCX file
            
        Returns:
            Path to DOCX file
        """
        try:
            from docx import Document
            from docx.shared import Pt
            
            # Read PDF
            reader = PdfReader(pdf_file)
            
            # Create Word document
            doc = Document()
            
            # Extract text from each page
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                
                # Add page number
                doc.add_heading(f'Page {page_num}', level=2)
                
                # Add text
                if text.strip():
                    paragraph = doc.add_paragraph(text)
                    paragraph.style.font.size = Pt(11)
                
                # Page break (except last page)
                if page_num < len(reader.pages):
                    doc.add_page_break()
            
            # Save document
            doc.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting PDF to Word: {str(e)}")
    
    @staticmethod
    def pdf_to_excel(pdf_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert PDF to Excel
        
        Note: This works best for PDFs with tabular data
        
        Args:
            pdf_file: PDF file object
            output_path: Path to save XLSX file
            
        Returns:
            Path to XLSX file
        """
        try:
            from openpyxl import Workbook
            
            # Read PDF
            reader = PdfReader(pdf_file)
            
            # Create workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "PDF Data"
            
            row_num = 1
            
            # Extract text from each page
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                
                # Add page header
                ws.cell(row=row_num, column=1, value=f"Page {page_num}")
                row_num += 1
                
                # Split text into lines and add to cells
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        # Try to split by common delimiters
                        cells = line.split('\t') if '\t' in line else [line]
                        for col_num, cell_value in enumerate(cells, 1):
                            ws.cell(row=row_num, column=col_num, value=cell_value.strip())
                        row_num += 1
                
                row_num += 1  # Blank row between pages
            
            # Save workbook
            wb.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting PDF to Excel: {str(e)}")
    
    @staticmethod
    def pdf_to_powerpoint(pdf_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert PDF to PowerPoint
        
        Args:
            pdf_file: PDF file object
            output_path: Path to save PPTX file
            
        Returns:
            Path to PPTX file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Read PDF
            reader = PdfReader(pdf_file)
            
            # Create presentation
            prs = Presentation()
            
            # Set slide size
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Extract text from each page and create slides
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                
                # Add slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
                
                # Set title
                title = slide.shapes.title
                title.text = f"Page {page_num}"
                
                # Add content
                if text.strip():
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.text = text[:500]  # Limit text length
            
            # Save presentation
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting PDF to PowerPoint: {str(e)}")
    
    @staticmethod
    def pdf_to_pdfa(pdf_file: BinaryIO, output_path: Path) -> Path:
        """
        Convert PDF to PDF/A (archival format)
        
        Args:
            pdf_file: PDF file object
            output_path: Path to save PDF/A file
            
        Returns:
            Path to PDF/A file
        """
        try:
            # Use pikepdf to convert to PDF/A
            pdf = pikepdf.open(pdf_file)
            
            # Add PDF/A metadata
            with pdf.open_metadata() as meta:
                meta['pdf:PDFVersion'] = '1.4'
                meta['pdfaid:part'] = '1'
                meta['pdfaid:conformance'] = 'B'
            
            # Save as PDF/A
            pdf.save(output_path, linearize=True)
            pdf.close()
            
            return output_path
            
        except Exception as e:
            raise Exception(f"Error converting to PDF/A: {str(e)}")
